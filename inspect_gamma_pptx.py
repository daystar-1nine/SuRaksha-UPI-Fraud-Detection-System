import pptx

prs = pptx.Presentation(r'C:\Users\suraj\Downloads\SURAKSHA-AI.pptx')
print(f"Total slides: {len(prs.slides)}")
print(f"Dimensions: {prs.slide_width.inches:.2f} x {prs.slide_height.inches:.2f}")

for idx, slide in enumerate(prs.slides):
    print(f"\n==================== SLIDE {idx+1} ====================")
    for s_idx, shape in enumerate(slide.shapes):
        shape_type = shape.shape_type
        left, top, w, h = shape.left.inches, shape.top.inches, shape.width.inches, shape.height.inches
        if shape.has_text_frame:
            texts = [p.text for p in shape.text_frame.paragraphs if p.text.strip()]
            joined_text = " // ".join(texts)
            print(f"  Shape {s_idx} [TEXT at ({left:.2f},{top:.2f}) {w:.2f}x{h:.2f}]: {joined_text}")
        elif shape.has_table:
            print(f"  Shape {s_idx} [TABLE at ({left:.2f},{top:.2f}) {w:.2f}x{h:.2f}]: {len(shape.table.rows)} rows x {len(shape.table.columns)} cols")
            for r in shape.table.rows:
                row_txt = [c.text.replace('\n', ' ') for c in r.cells]
                print(f"     Row: {' | '.join(row_txt)}")
        elif shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            print(f"  Shape {s_idx} [PICTURE at ({left:.2f},{top:.2f}) {w:.2f}x{h:.2f}]")
        else:
            print(f"  Shape {s_idx} [SHAPE {shape_type} at ({left:.2f},{top:.2f}) {w:.2f}x{h:.2f}]")
