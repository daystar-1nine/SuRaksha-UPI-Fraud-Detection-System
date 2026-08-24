import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # completely blank

    # Color Palette - Premium Fintech & Cybersecurity
    BG_DARK = RGBColor(7, 10, 19)        # #070a13
    CARD_BG = RGBColor(15, 23, 42)       # #0f172a
    CARD_BORDER = RGBColor(30, 41, 59)   # #1e293b
    CARD_HIGHLIGHT = RGBColor(30, 58, 138) # #1e3a8a
    
    TEXT_WHITE = RGBColor(255, 255, 255)
    TEXT_MUTED = RGBColor(148, 163, 184) # #94a3b8
    TEXT_DIM = RGBColor(100, 116, 139)   # #64748b
    
    CYAN_ACCENT = RGBColor(56, 189, 248)  # #38bdf8
    BLUE_ACCENT = RGBColor(59, 130, 246)  # #3b82f6
    GREEN_SAFE = RGBColor(16, 185, 129)   # #10b981
    RED_ALERT = RGBColor(239, 68, 68)     # #ef4444
    AMBER_WARN = RGBColor(245, 158, 11)   # #f59e0b

    assets_dir = r"S:\Hackathon\SuRaksha\assets\screenshots"

    def apply_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_DARK
        bg.line.fill.background()
        return bg

    def add_header(slide, tag_text, title_text, category_color=CYAN_ACCENT):
        # Category Tag
        tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf = tag_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = tag_text.upper()
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = category_color

        # Main Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.7), Inches(0.7))
        tf2 = title_box.text_frame
        tf2.word_wrap = True
        tf2.margin_left = tf2.margin_top = tf2.margin_right = tf2.margin_bottom = 0
        p2 = tf2.paragraphs[0]
        p2.text = title_text
        p2.font.size = Pt(24)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_WHITE

    def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=CARD_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
        return card

    # ==========================================
    # SLIDE 1: COVER / HOOK
    # ==========================================
    s1 = prs.slides.add_slide(blank_layout)
    apply_bg(s1)
    
    # Left Hero Container
    add_card(s1, 0.8, 0.8, 7.2, 5.9, CARD_BG, CARD_HIGHLIGHT)
    
    # Competition Badge
    badge = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.2), Inches(5.8), Inches(0.45))
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(30, 58, 138)
    badge.line.color.rgb = CYAN_ACCENT
    badge.line.width = Pt(1)
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = "EUREKA! PITCH COMPETITION 2026 | SJCEM × IIT BOMBAY NEC"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT
    p.alignment = PP_ALIGN.CENTER

    # Brand Title
    tb = s1.shapes.add_textbox(Inches(1.2), Inches(1.85), Inches(6.5), Inches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "SURAKSHA AI"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    
    p2 = tf.add_paragraph()
    p2.text = "Real-Time Pre-Payment Risk & Fraud Intelligence"
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.color.rgb = CYAN_ACCENT

    # Slogan / Value Proposition
    tb_val = s1.shapes.add_textbox(Inches(1.2), Inches(3.6), Inches(6.4), Inches(1.5))
    tf_val = tb_val.text_frame
    tf_val.word_wrap = True
    p = tf_val.paragraphs[0]
    p.text = '"The payment is instant. The warning should be too."'
    p.font.size = Pt(15)
    p.font.italic = True
    p.font.color.rgb = RGBColor(226, 232, 240)
    
    p2 = tf_val.add_paragraph()
    p2.text = "An intelligent, privacy-first security layer that evaluates QR codes, UPI IDs, payment screenshots, and phishing links before money leaves your account."
    p2.font.size = Pt(13)
    p2.font.color.rgb = TEXT_MUTED

    # Highlights Pills
    pills = [
        ("⚡ <200ms Edge Latency", GREEN_SAFE),
        ("🛡️ Multi-Modal AI Shield", BLUE_ACCENT),
        ("🔒 Zero-Data Retention", CYAN_ACCENT)
    ]
    for i, (txt, col) in enumerate(pills):
        pill = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2 + i*2.15), Inches(5.8), Inches(2.05), Inches(0.45))
        pill.fill.solid()
        pill.fill.fore_color.rgb = RGBColor(15, 23, 42)
        pill.line.color.rgb = col
        pill.line.width = Pt(1)
        tf = pill.text_frame
        p = tf.paragraphs[0]
        p.text = txt
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = col
        p.alignment = PP_ALIGN.CENTER

    # Right Hero Image Card (Logo & UI Graphic)
    add_card(s1, 8.3, 0.8, 4.2, 5.9, CARD_BG, CARD_BORDER)
    logo_path = os.path.join(assets_dir, "real_life_usecase.png")
    if os.path.exists(logo_path):
        s1.shapes.add_picture(logo_path, Inches(8.5), Inches(1.1), width=Inches(3.8))
        
    caption_box = s1.shapes.add_textbox(Inches(8.5), Inches(5.7), Inches(3.8), Inches(0.8))
    tf_c = caption_box.text_frame
    tf_c.word_wrap = True
    p = tf_c.paragraphs[0]
    p.text = "Tested on 4,300+ real-world payment attack scenarios"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = GREEN_SAFE
    p.alignment = PP_ALIGN.CENTER

    # ==========================================
    # SLIDE 2: THE PROBLEM
    # ==========================================
    s2 = prs.slides.add_slide(blank_layout)
    apply_bg(s2)
    add_header(s2, "THE DIGITAL PAYMENT CRISIS", "Instant Payments. Zero Pre-Transaction Defense.", RED_ALERT)

    # 3 Problem Cards
    problems = [
        ("01", "QR Sticker Swapping", "Physical Retail Tampering", 
         "Fraudsters paste deceptive scam QR stickers over legitimate counter boards in retail shops. The transaction completes instantly to the attacker's account.",
         "₹100+ Crore annual loss across retail UPI", RED_ALERT),
        ("02", "Doctored Payment Screenshots", "Digital Visual Forgery", 
         "Customers flash fabricated GPay / PhonePe receipt screenshots made via image editors or fake receipt apps. Merchants hand over goods without receiving actual funds.",
         "Merchants lose daily working capital", AMBER_WARN),
        ("03", "Phishing & Urgency Traps", "Social Engineering Traps", 
         "Scammers send WhatsApp / SMS links disguised as 'Cashback' or 'KYC Update'. Clicking triggers a debit payment request (`upi://pay`) instead of a credit.",
         "80%+ victims realize only after debit", RED_ALERT)
    ]

    for i, (num, title, subtitle, desc, stat, color) in enumerate(problems):
        x = 0.8 + i * 4.0
        add_card(s2, x, 1.6, 3.7, 5.1, CARD_BG, CARD_BORDER)
        
        # Top Accent Line
        acc = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.6), Inches(3.7), Inches(0.08))
        acc.fill.solid()
        acc.fill.fore_color.rgb = color
        acc.line.fill.background()

        # Number Tag
        num_box = s2.shapes.add_textbox(Inches(x + 0.3), Inches(1.9), Inches(1.0), Inches(0.6))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = color

        # Content
        c_box = s2.shapes.add_textbox(Inches(x + 0.3), Inches(2.6), Inches(3.1), Inches(3.8))
        tf = c_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE

        p_sub = tf.add_paragraph()
        p_sub.text = subtitle
        p_sub.font.size = Pt(11)
        p_sub.font.color.rgb = CYAN_ACCENT

        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = TEXT_MUTED

        p_stat = tf.add_paragraph()
        p_stat.text = "⚠️ " + stat
        p_stat.font.size = Pt(11)
        p_stat.font.bold = True
        p_stat.font.color.rgb = color

    # ==========================================
    # SLIDE 3: THE CORE INSIGHT
    # ==========================================
    s3 = prs.slides.add_slide(blank_layout)
    apply_bg(s3)
    add_header(s3, "THE CORE INSIGHT", "The Security Gap Exists BEFORE The PIN Is Entered", BLUE_ACCENT)

    # Comparison Grid: Traditional vs Suraksha
    add_card(s3, 0.8, 1.6, 5.7, 5.1, CARD_BG, RGBColor(75, 85, 99))
    add_card(s3, 6.8, 1.6, 5.7, 5.1, CARD_BG, RGBColor(30, 58, 138))

    # Left: Traditional Approach
    tb_left = s3.shapes.add_textbox(Inches(1.2), Inches(1.9), Inches(4.9), Inches(4.5))
    tf_l = tb_left.text_frame
    tf_l.word_wrap = True
    
    p = tf_l.paragraphs[0]
    p.text = "❌ Current Digital Payment Apps"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RED_ALERT

    points_l = [
        ("Reactive Focus: ", "Fraud is flagged only AFTER money is deducted via post-transaction bank dispute forms."),
        ("No QR Verification: ", "Apps blindly parse the raw UPI URL without validating merchant sticker authenticity."),
        ("No Receipt Verification: ", "Merchants have zero built-in tools to detect Photoshop-edited receipts."),
        ("English-Centric: ", "Complex bank terms confuse rural, non-English speaking retail merchants.")
    ]
    for bold_txt, norm_txt in points_l:
        p = tf_l.add_paragraph()
        p.text = "• " + bold_txt + norm_txt
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_MUTED

    # Right: Suraksha AI Approach
    tb_right = s3.shapes.add_textbox(Inches(7.2), Inches(1.9), Inches(4.9), Inches(4.5))
    tf_r = tb_right.text_frame
    tf_r.word_wrap = True
    
    p = tf_r.paragraphs[0]
    p.text = "🛡️ Suraksha AI: Zero-Trust Pre-Payment Shield"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN_SAFE

    points_r = [
        ("Pre-Authorization Defense: ", "Analyzes threat parameters in <200ms BEFORE the user types their 4/6-digit PIN."),
        ("Cryptographic QR Shield: ", "Validates HMAC SHA-256 merchant signatures to instantly expose sticker swaps."),
        ("In-Memory ELA Forensics: ", "Applies 75% Q-factor Error Level Analysis in RAM to expose image pixel tampering."),
        ("Bilingual Intelligence: ", "Seamless English & Devanagari Hindi translation protecting Tier-2/3 users.")
    ]
    for bold_txt, norm_txt in points_r:
        p = tf_r.add_paragraph()
        p.text = "✓ " + bold_txt + norm_txt
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_WHITE

    # ==========================================
    # SLIDE 4: THE SOLUTION (SURAKSHA AI)
    # ==========================================
    s4 = prs.slides.add_slide(blank_layout)
    apply_bg(s4)
    add_header(s4, "THE SOLUTION", "Suraksha AI: Multi-Modal Pre-Payment Defense Ecosystem", CYAN_ACCENT)

    # 4 Architecture Pillars
    pillars = [
        ("01. Live QR Scanner", "WebRTC camera feed parses UPI URI, checks VPA blacklists & validates HMAC signatures.", GREEN_SAFE),
        ("02. Image Forensic ELA", "RAM-only OpenCV JPEG Error Level Analysis detects doctored receipt screenshots.", BLUE_ACCENT),
        ("03. NLP Message Validator", "Naive Bayes & Regex heuristic matrix flags cashback / KYC urgency traps in <20ms.", AMBER_WARN),
        ("04. Signed QR Generator", "Empowers shopkeepers to generate cryptographically tamper-proof merchant QR codes.", CYAN_ACCENT)
    ]

    for i, (title, desc, col) in enumerate(pillars):
        x = 0.8 + i * 2.95
        add_card(s4, x, 1.6, 2.75, 2.3, CARD_BG, CARD_BORDER)
        
        tb = s4.shapes.add_textbox(Inches(x + 0.2), Inches(1.8), Inches(2.35), Inches(1.9))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = col
        
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = TEXT_MUTED

    # Bottom Banner: How Decisions Are Made
    add_card(s4, 0.8, 4.15, 11.7, 2.55, RGBColor(15, 23, 42), RGBColor(56, 189, 248))
    
    tb_flow = s4.shapes.add_textbox(Inches(1.1), Inches(4.3), Inches(11.1), Inches(2.2))
    tf_f = tb_flow.text_frame
    tf_f.word_wrap = True
    
    p = tf_f.paragraphs[0]
    p.text = "THE ZERO-TRUST EVALUATION PIPELINE"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT

    p_steps = tf_f.add_paragraph()
    p_steps.text = "User Input (QR / Image / SMS)  ──>  In-Memory Sanitization  ──>  Multi-Signal Evaluation (ELA + HMAC + NLP)  ──>  Weighted Risk Engine  ──>  Actionable HUD Verdict"
    p_steps.font.size = Pt(13)
    p_steps.font.bold = True
    p_steps.font.color.rgb = TEXT_WHITE

    p_scores = tf_f.add_paragraph()
    p_scores.text = "🟢 Score 0–30: SAFE (Launch Payment App)   |   🟡 Score 31–70: CAUTION (Verify Merchant Identity)   |   🔴 Score 71–100: HIGH RISK (Block Transaction & Report)"
    p_scores.font.size = Pt(12)
    p_scores.font.color.rgb = RGBColor(226, 232, 240)

    # ==========================================
    # SLIDE 5: LIVE PRODUCT DEMO (EVIDENCE)
    # ==========================================
    s5 = prs.slides.add_slide(blank_layout)
    apply_bg(s5)
    add_header(s5, "PRODUCT EVIDENCE", "Live Working Product: Real-Time Threat Interception", GREEN_SAFE)

    # 3 Screenshots from real product
    demo_cards = [
        ("1. Multi-Modal Scanner Hub", "realtime_scanner.png", "Live camera QR & screenshot upload interface"),
        ("2. Explainable Threat HUD", "threat_breakdown_hud.png", "Calculates composite risk score with factor breakdown"),
        ("3. Forensic ELA Analysis", "forensic_ela_check.png", "Exposes pixel modifications in doctored receipts")
    ]

    for i, (title, img_name, sub) in enumerate(demo_cards):
        x = 0.8 + i * 4.0
        add_card(s5, x, 1.6, 3.7, 5.1, CARD_BG, CARD_BORDER)
        
        # Title
        tb = s5.shapes.add_textbox(Inches(x + 0.2), Inches(1.75), Inches(3.3), Inches(0.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = CYAN_ACCENT
        
        # Image
        img_path = os.path.join(assets_dir, img_name)
        if os.path.exists(img_path):
            s5.shapes.add_picture(img_path, Inches(x + 0.2), Inches(2.3), width=Inches(3.3))
            
        # Subtitle
        tb_sub = s5.shapes.add_textbox(Inches(x + 0.2), Inches(6.0), Inches(3.3), Inches(0.6))
        tf_s = tb_sub.text_frame
        tf_s.word_wrap = True
        p = tf_s.paragraphs[0]
        p.text = sub
        p.font.size = Pt(10.5)
        p.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 6: HOW IT WORKS (TECH & INNOVATION)
    # ==========================================
    s6 = prs.slides.add_slide(blank_layout)
    apply_bg(s6)
    add_header(s6, "TECHNICAL ARCHITECTURE", "High-Performance, Privacy-Preserving Engine", BLUE_ACCENT)

    # 3 Column Tech Breakdown
    tech_cols = [
        ("Layer 1: Visual Image Forensics", [
            ("75% JPEG Error Level Analysis: ", "Re-compresses in-memory byte streams to detect localized quantization noise."),
            ("Laplacian Spatial Variance: ", "Measures edge sharpness inconsistencies to detect sharp text overlays (Threshold > 2000)."),
            ("Zero-Disk Persistence: ", "Bytes processed via `io.BytesIO` in volatile RAM without saving files to disk.")
        ], CYAN_ACCENT),
        
        ("Layer 2: Cryptographic QR Shield", [
            ("HMAC SHA-256 Signatures: ", "Merchants sign VPA parameters using private keys (`sign = SHA256(pn || pa || key)`)."),
            ("Merchant Registry Lookup: ", "Matches scanned signatures against trusted databases to expose sticker swapping."),
            ("Format & Scheme Integrity: ", "Validates strictly formatted `upi://pay` URIs against spoofed deep links.")
        ], BLUE_ACCENT),
        
        ("Layer 3: NLP & Threat Intelligence", [
            ("Naive Bayes TF-IDF Model: ", "Classifies SMS/chat strings into benign alerts vs social engineering traps in <20ms."),
            ("Regex Heuristic Rules: ", "Detects high-pressure triggers (account blocked, KYC update, cashback reward)."),
            ("Bilingual Translation: ", "Regex word-boundary mapping ensures zero Hindi UI distortion.")
        ], GREEN_SAFE)
    ]

    for i, (title, bullets, col) in enumerate(tech_cols):
        x = 0.8 + i * 4.0
        add_card(s6, x, 1.6, 3.7, 5.1, CARD_BG, CARD_BORDER)
        
        # Header Box
        tb = s6.shapes.add_textbox(Inches(x + 0.25), Inches(1.8), Inches(3.2), Inches(0.7))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = col

        # Bullet content
        tb_c = s6.shapes.add_textbox(Inches(x + 0.25), Inches(2.6), Inches(3.2), Inches(3.9))
        tf_c = tb_c.text_frame
        tf_c.word_wrap = True
        for b_title, b_desc in bullets:
            p = tf_c.add_paragraph() if tf_c.paragraphs[0].text else tf_c.paragraphs[0]
            p.text = "• " + b_title + b_desc
            p.font.size = Pt(11.5)
            p.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 7: WHY SURAKSHA? (DIFFERENTIATION)
    # ==========================================
    s7 = prs.slides.add_slide(blank_layout)
    apply_bg(s7)
    add_header(s7, "COMPETITIVE DIFFERENTIATION", "From Reactive Disputes to Pre-Payment Protection", CYAN_ACCENT)

    # Comparison Table
    add_card(s7, 0.8, 1.6, 11.7, 5.1, CARD_BG, CARD_BORDER)
    
    table_shape = s7.shapes.add_table(6, 4, Inches(1.1), Inches(1.9), Inches(11.1), Inches(4.4))
    table = table_shape.table
    table.columns[0].width = Inches(3.3)
    table.columns[1].width = Inches(2.6)
    table.columns[2].width = Inches(2.6)
    table.columns[3].width = Inches(2.6)

    headers = ["Capability / Feature", "Traditional UPI Apps", "Server-Side Fraud Logs", "Suraksha AI Platform"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = CYAN_ACCENT if j == 3 else TEXT_WHITE
        p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT

    rows = [
        ("Inspection Timing", "Post-PIN Entry", "Post-Settlement Logs", "Pre-Authorization (<200ms)"),
        ("QR Sticker Swap Defense", "None (Blind trust)", "None", "HMAC SHA-256 Validation"),
        ("Screenshot Forgery Detection", "None", "None", "In-Memory ELA Forensics"),
        ("User Privacy Policy", "Server Transaction Logs", "Bulk Server Storage", "Zero-Data Retention (RAM)"),
        ("Tier-2/3 Language Inclusivity", "Static / English", "English Only", "Dynamic English & Hindi")
    ]

    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i+1, j)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            if j == 3:
                p.font.bold = True
                p.font.color.rgb = GREEN_SAFE
            elif j == 0:
                p.font.bold = True
                p.font.color.rgb = TEXT_WHITE
            else:
                p.font.color.rgb = TEXT_MUTED
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT

    # ==========================================
    # SLIDE 8: BUSINESS MODEL & SCALE
    # ==========================================
    s8 = prs.slides.add_slide(blank_layout)
    apply_bg(s8)
    add_header(s8, "BUSINESS MODEL & SCALE", "Built for Users. Scaled Through Ecosystem APIs.", GREEN_SAFE)

    # 3 Revenue & Go-to-Market Cards
    biz_pillars = [
        ("B2C: Consumer Shield", "Freemium Mobile App", [
            ("Core Scanner: ", "Free real-time QR & SMS phishing protection for everyday users."),
            ("Premium Security HUD: ", "Subscription tier with family protection, instant scam reporting, and cyber insurance link.")
        ], GREEN_SAFE),
        
        ("B2B: Merchant Protection", "Retail Merchant Subscriptions", [
            ("Signed QR Certification: ", "Monthly subscription for verified shopkeeper badges and tamper-proof counter stands."),
            ("POS Screenshot Verifier: ", "Instant audio chime and visual screen verification to protect shopkeepers from fake receipts.")
        ], BLUE_ACCENT),
        
        ("B2B2C / API: FinTech SDK", "Enterprise Risk-as-a-Service", [
            ("Pre-Auth Risk API: ", "Microservice API licensing for neo-banks, payment aggregators, and digital wallets."),
            ("Fraud Intelligence Feeds: ", "Shared scammer VPA blacklist telemetry for financial institutions.")
        ], CYAN_ACCENT)
    ]

    for i, (title, subtitle, bullets, col) in enumerate(biz_pillars):
        x = 0.8 + i * 4.0
        add_card(s8, x, 1.6, 3.7, 5.1, CARD_BG, CARD_BORDER)
        
        # Header
        tb = s8.shapes.add_textbox(Inches(x + 0.25), Inches(1.8), Inches(3.2), Inches(0.9))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = col
        
        p_sub = tf.add_paragraph()
        p_sub.text = subtitle
        p_sub.font.size = Pt(11)
        p_sub.font.color.rgb = TEXT_WHITE

        # Bullets
        tb_c = s8.shapes.add_textbox(Inches(x + 0.25), Inches(2.8), Inches(3.2), Inches(3.7))
        tf_c = tb_c.text_frame
        tf_c.word_wrap = True
        for b_title, b_desc in bullets:
            p = tf_c.add_paragraph() if tf_c.paragraphs[0].text else tf_c.paragraphs[0]
            p.text = "• " + b_title + b_desc
            p.font.size = Pt(11.5)
            p.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 9: VISION & CLOSING
    # ==========================================
    s9 = prs.slides.add_slide(blank_layout)
    apply_bg(s9)

    # Large Central Card
    add_card(s9, 1.2, 1.0, 10.9, 5.5, CARD_BG, RGBColor(30, 58, 138))

    # Badge
    b_box = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.6), Inches(1.4), Inches(4.1), Inches(0.45))
    b_box.fill.solid()
    b_box.fill.fore_color.rgb = RGBColor(15, 23, 42)
    b_box.line.color.rgb = CYAN_ACCENT
    b_box.line.width = Pt(1)
    tf = b_box.text_frame
    p = tf.paragraphs[0]
    p.text = "THE SURAKSHA AI VISION"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT
    p.alignment = PP_ALIGN.CENTER

    # Main Vision Statement
    tb_vis = s9.shapes.add_textbox(Inches(1.6), Inches(2.1), Inches(10.1), Inches(2.2))
    tf_v = tb_vis.text_frame
    tf_v.word_wrap = True
    p = tf_v.paragraphs[0]
    p.text = "Don't discover the scam after you pay.\nDiscover the risk before you pay."
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.alignment = PP_ALIGN.CENTER

    p_sub = tf_v.add_paragraph()
    p_sub.text = "Digital payments should not only be fast — they must be informed, transparent, and secure."
    p_sub.font.size = Pt(16)
    p_sub.font.color.rgb = CYAN_ACCENT
    p_sub.alignment = PP_ALIGN.CENTER

    # Execution Milestones
    m_boxes = [
        ("STAGE 1: PROTOTYPE", "Working PWA, Flask microservice & OpenCV ELA tested on 4,300 samples.", GREEN_SAFE),
        ("STAGE 2: VALIDATION", "Pilot trials with local retail shopkeepers in Tier-2/3 market corridors.", BLUE_ACCENT),
        ("STAGE 3: SCALE", "FinTech API licensing & multi-bank threat intelligence sharing.", CYAN_ACCENT)
    ]

    for i, (m_title, m_desc, col) in enumerate(m_boxes):
        x = 1.6 + i * 3.45
        card_m = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(4.5), Inches(3.2), Inches(1.5))
        card_m.fill.solid()
        card_m.fill.fore_color.rgb = RGBColor(15, 23, 42)
        card_m.line.color.rgb = col
        card_m.line.width = Pt(1)
        
        tf_m = card_m.text_frame
        tf_m.word_wrap = True
        p = tf_m.paragraphs[0]
        p.text = m_title
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = col
        
        p_desc = tf_m.add_paragraph()
        p_desc.text = m_desc
        p_desc.font.size = Pt(10)
        p_desc.font.color.rgb = TEXT_MUTED

    # Save to disk
    output_path = r"S:\Hackathon\SuRaksha\Suraksha_AI_Eureka_Pitch_2026.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    build_presentation()
