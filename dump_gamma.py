import pptx

prs = pptx.Presentation(r'C:\Users\suraj\Downloads\SURAKSHA-AI.pptx')
with open(r'S:\Hackathon\SuRaksha\gamma_structure.txt', 'w', encoding='utf-8') as f:
    f.write(f"Total slides: {len(prs.slides)}\n")
    f.write(f"Dimensions: {prs.slide_width.inches:.2f} x {prs.slide_height.inches:.2f}\n")

    for idx, slide in enumerate(prs.slides):
        f.write(f"\n==================== SLIDE {idx+1} ====================\n")
        for s_idx, shape in enumerate(slide.shapes):
            shape_type = shape.shape_type
            left, top, w, h = shape.left.inches, shape.top.inches, shape.width.inches, shape.height.inches
            if shape.has_text_frame:
                texts = [p.text for p in shape.text_frame.paragraphs if p.text.strip()]
                joined_text = " // ".join(texts)
                f.write(f"  Shape {s_idx} [TEXT at ({left:.2f},{top:.2f}) {w:.2f}x{h:.2f}]: {joined_text}\n")
            elif shape.has_table:
                f.write(f"  Shape {s_idx} [TABLE at ({left:.2f},{top:.2f}) {w:.2f}x{h:.2f}]: {len(shape.table.rows)} rows x {len(shape.table.columns)} cols\n")
                for r in shape.table.rows:
                    row_txt = [c.text.replace('\n', ' ') for c in r.cells]
                    f.write(f"     Row: {' | '.join(row_txt)}\n")
            elif shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
                f.write(f"  Shape {s_idx} [PICTURE at ({left:.2f},{top:.2f}) {w:.2f}x{h:.2f}]\n")
            else:
                f.write(f"  Shape {s_idx} [SHAPE {shape_type} at ({left:.2f},{top:.2f}) {w:.2f}x{h:.2f}]\n")

print("Wrote structure to gamma_structure.txt")
