import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def build_enhanced_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Premium Fintech & Cybersecurity Color Palette
    BG_DARK = RGBColor(7, 10, 19)          # #070a13 Deep Space Navy
    CARD_BG = RGBColor(15, 23, 42)         # #0f172a Deep Slate
    CARD_BORDER = RGBColor(30, 41, 59)     # #1e293b Subtle Slate Border
    CARD_ACCENT_BLUE = RGBColor(30, 58, 138) # #1e3a8a Active Highlight
    
    TEXT_WHITE = RGBColor(255, 255, 255)
    TEXT_LIGHT = RGBColor(241, 245, 249)   # #f1f5f9 Crisp
    TEXT_MUTED = RGBColor(148, 163, 184)   # #94a3b8 Slate Muted
    TEXT_DIM = RGBColor(100, 116, 139)     # #64748b
    
    CYAN_ACCENT = RGBColor(56, 189, 248)    # #38bdf8 Electric Cyan
    BLUE_ACCENT = RGBColor(59, 130, 246)    # #3b82f6 Cobalt Blue
    GREEN_SAFE = RGBColor(16, 185, 129)     # #10b981 Emerald Safe
    RED_ALERT = RGBColor(239, 68, 68)       # #ef4444 Crimson Alert
    AMBER_WARN = RGBColor(245, 158, 11)     # #f59e0b Amber Caution

    assets_dir = r"S:\Hackathon\SuRaksha\assets\screenshots"
    logo_file = os.path.join(assets_dir, "logo.png")

    def apply_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_DARK
        bg.line.fill.background()
        return bg

    def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=CARD_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.2)
        return card

    def add_header(slide, tag_text, title_text, category_color=CYAN_ACCENT):
        # Header Badge / Tag
        tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.38), Inches(10.0), Inches(0.32))
        tf = tag_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = tag_text.upper()
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = category_color

        # Main Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.68), Inches(10.5), Inches(0.65))
        tf2 = title_box.text_frame
        tf2.word_wrap = True
        tf2.margin_left = tf2.margin_top = tf2.margin_right = tf2.margin_bottom = 0
        p2 = tf2.paragraphs[0]
        p2.text = title_text
        p2.font.size = Pt(23)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_WHITE

        # Logo on Top Right if exists
        if os.path.exists(logo_file):
            try:
                slide.shapes.add_picture(logo_file, Inches(11.8), Inches(0.4), height=Inches(0.65))
            except Exception:
                pass

    # ==========================================
    # SLIDE 1: COVER / THE HOOK
    # ==========================================
    s1 = prs.slides.add_slide(blank_layout)
    apply_bg(s1)
    
    # Left Hero Container
    add_card(s1, 0.8, 0.8, 7.3, 5.9, CARD_BG, CARD_ACCENT_BLUE)
    
    # Competition Header Badge
    badge = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.15), Inches(5.6), Inches(0.42))
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(30, 58, 138)
    badge.line.color.rgb = CYAN_ACCENT
    badge.line.width = Pt(1)
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = "EUREKA! PITCH COMPETITION 2026 | SJCEM × IIT BOMBAY NEC"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT
    p.alignment = PP_ALIGN.CENTER

    # Brand Title
    tb = s1.shapes.add_textbox(Inches(1.2), Inches(1.75), Inches(6.5), Inches(1.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "SURAKSHA AI"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    
    p2 = tf.add_paragraph()
    p2.text = "Real-Time Payment Risk & Fraud Intelligence Platform"
    p2.font.size = Pt(17)
    p2.font.bold = True
    p2.font.color.rgb = CYAN_ACCENT

    # Slogan / Pitch Hook
    tb_val = s1.shapes.add_textbox(Inches(1.2), Inches(3.45), Inches(6.5), Inches(1.5))
    tf_val = tb_val.text_frame
    tf_val.word_wrap = True
    p = tf_val.paragraphs[0]
    p.text = '"The payment is instant. The warning should be too."'
    p.font.size = Pt(15)
    p.font.italic = True
    p.font.color.rgb = RGBColor(226, 232, 240)
    
    p2 = tf_val.add_paragraph()
    p2.text = "An intelligent, privacy-first security layer evaluating QR codes, UPI IDs, payment screenshots, and phishing links before money leaves your account."
    p2.font.size = Pt(12.5)
    p2.font.color.rgb = TEXT_MUTED

    # Highlights Pills (Clean, no emojis)
    pills = [
        ("[FAST] <200ms Execution", GREEN_SAFE),
        ("[SHIELD] Multi-Modal AI", BLUE_ACCENT),
        ("[PRIVACY] Zero-Data Retention", CYAN_ACCENT)
    ]
    for i, (txt, col) in enumerate(pills):
        pill = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2 + i*2.15), Inches(5.8), Inches(2.05), Inches(0.42))
        pill.fill.solid()
        pill.fill.fore_color.rgb = RGBColor(15, 23, 42)
        pill.line.color.rgb = col
        pill.line.width = Pt(1)
        tf = pill.text_frame
        p = tf.paragraphs[0]
        p.text = txt
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = col
        p.alignment = PP_ALIGN.CENTER

    # Right Hero Image Card
    add_card(s1, 8.4, 0.8, 4.1, 5.9, CARD_BG, CARD_BORDER)
    hero_img = os.path.join(assets_dir, "real_life_usecase.png")
    if os.path.exists(hero_img):
        s1.shapes.add_picture(hero_img, Inches(8.55), Inches(1.05), width=Inches(3.8))
        
    caption_box = s1.shapes.add_textbox(Inches(8.55), Inches(5.75), Inches(3.8), Inches(0.7))
    tf_c = caption_box.text_frame
    tf_c.word_wrap = True
    p = tf_c.paragraphs[0]
    p.text = "[ACTIVE PROTOTYPE] Validated across real payment attack scenarios"
    p.font.size = Pt(10.5)
    p.font.bold = True
    p.font.color.rgb = GREEN_SAFE
    p.alignment = PP_ALIGN.CENTER

    s1.notes_slide.notes_text_frame.text = (
        "[TIME: 0:00 - 0:30]\n"
        "[WHAT TO SAY]: Judges, imagine you walk up to a busy retail counter. You scan the QR code, type ₹500, enter your UPI PIN, and hear the success chime. But two minutes later, the shopkeeper tells you he never received the money. Why? Because a scammer pasted a duplicate QR sticker right over the shop's original board. Or consider a small vendor who hands over ₹5,000 worth of goods because a customer flashes a doctored Google Pay confirmation screenshot. India processes billions of instant payments monthly, but fraud awareness takes hours. That is why we built Suraksha AI.\n"
        "[ACTION]: Stand center stage, confident posture, no slide notes in hand.\n"
        "[TRANSITION]: Let's look at the three critical fraud vectors plaguing our digital economy."
    )

    # ==========================================
    # SLIDE 2: THE PROBLEM
    # ==========================================
    s2 = prs.slides.add_slide(blank_layout)
    apply_bg(s2)
    add_header(s2, "THE PROBLEM", "Instant Payments. Delayed Fraud Awareness.", RED_ALERT)

    problems = [
        ("01", "Physical QR Sticker Swapping", "Retail Point-of-Sale Tampering", 
         "Fraudsters paste deceptive scam QR stickers over legitimate counter boards in retail shops. Funds are silently diverted to attacker accounts at checkout.",
         "Threat: Silent checkout fund diversion", RED_ALERT),
        ("02", "Doctored Payment Screenshots", "Digital Visual Image Forgery", 
         "Customers flash fabricated GPay / PhonePe receipt screenshots made via photo editing apps. Merchants hand over goods without receiving actual funds.",
         "Threat: Merchants lose daily working capital", AMBER_WARN),
        ("03", "Phishing & Urgency Traps", "Social Engineering Payment Links", 
         "Scammers send WhatsApp / SMS links disguised as 'Cashback' or 'KYC Update'. Clicking triggers a debit payment request (`upi://pay`) instead of a credit.",
         "Threat: Victims realize only after bank debit", RED_ALERT)
    ]

    for i, (num, title, subtitle, desc, stat, color) in enumerate(problems):
        x = 0.8 + i * 4.0
        add_card(s2, x, 1.55, 3.7, 5.25, CARD_BG, CARD_BORDER)
        
        acc = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.55), Inches(3.7), Inches(0.08))
        acc.fill.solid()
        acc.fill.fore_color.rgb = color
        acc.line.fill.background()

        num_box = s2.shapes.add_textbox(Inches(x + 0.25), Inches(1.75), Inches(1.0), Inches(0.55))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = color

        c_box = s2.shapes.add_textbox(Inches(x + 0.25), Inches(2.4), Inches(3.2), Inches(4.1))
        tf = c_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(15.5)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE

        p_sub = tf.add_paragraph()
        p_sub.text = subtitle
        p_sub.font.size = Pt(11)
        p_sub.font.color.rgb = CYAN_ACCENT

        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(11.5)
        p_desc.font.color.rgb = TEXT_MUTED

        p_stat = tf.add_paragraph()
        p_stat.text = "[VULNERABILITY] " + stat
        p_stat.font.size = Pt(10.5)
        p_stat.font.bold = True
        p_stat.font.color.rgb = color

    s2.notes_slide.notes_text_frame.text = (
        "[TIME: 0:30 - 1:10]\n"
        "[WHAT TO SAY]: Digital payments are fast and seamless, but fraud tactics have evolved to exploit the human interface before bank settlement. First, physical QR sticker swapping, where retail merchant boards are silently hijacked. Second, doctored confirmation screenshots, where fake receipts dupe vendors out of inventory. And third, social engineering phishing links that disguise payment requests as cashback rewards. The core problem is that users lack real-time risk intelligence at the exact moment of payment.\n"
        "[ACTION]: Point to each of the 3 problem cards.\n"
        "[TRANSITION]: Why does this vulnerability exist in our current apps? That brings us to our core insight."
    )

    # ==========================================
    # SLIDE 3: THE INSIGHT
    # ==========================================
    s3 = prs.slides.add_slide(blank_layout)
    apply_bg(s3)
    add_header(s3, "THE CORE INSIGHT", "The Security Gap Exists Before The Payment", BLUE_ACCENT)

    add_card(s3, 0.8, 1.55, 5.7, 5.25, CARD_BG, RGBColor(75, 85, 99))
    add_card(s3, 6.8, 1.55, 5.7, 5.25, CARD_BG, RGBColor(30, 58, 138))

    tb_left = s3.shapes.add_textbox(Inches(1.15), Inches(1.8), Inches(5.0), Inches(4.7))
    tf_l = tb_left.text_frame
    tf_l.word_wrap = True
    
    p = tf_l.paragraphs[0]
    p.text = "[TRADITIONAL] Current Payment Apps"
    p.font.size = Pt(16.5)
    p.font.bold = True
    p.font.color.rgb = RED_ALERT

    points_l = [
        ("Reactive Focus: ", "Fraud is flagged only AFTER money is deducted via post-transaction dispute forms."),
        ("Blind Trust in QR URLs: ", "Apps parse the raw link without checking physical sticker validity."),
        ("No Receipt Verification: ", "Merchants have zero tools to detect Photoshop-edited receipts."),
        ("Language Barrier: ", "Technical security jargon confuses non-English speaking vendors.")
    ]
    for bold_txt, norm_txt in points_l:
        p = tf_l.add_paragraph()
        p.text = "• " + bold_txt + norm_txt
        p.font.size = Pt(11.5)
        p.font.color.rgb = TEXT_MUTED

    tb_right = s3.shapes.add_textbox(Inches(7.15), Inches(1.8), Inches(5.0), Inches(4.7))
    tf_r = tb_right.text_frame
    tf_r.word_wrap = True
    
    p = tf_r.paragraphs[0]
    p.text = "[PRE-AUTH SHIELD] Suraksha AI"
    p.font.size = Pt(16.5)
    p.font.bold = True
    p.font.color.rgb = GREEN_SAFE

    points_r = [
        ("Pre-Payment Warning: ", "Analyzes threat parameters in <200ms BEFORE the user types their PIN."),
        ("Cryptographic QR Shield: ", "Validates HMAC SHA-256 merchant signatures to catch sticker swaps."),
        ("In-Memory ELA Forensics: ", "Applies 75% Q-factor Error Level Analysis in RAM to expose image tampering."),
        ("Bilingual Intelligence: ", "Seamless English & Devanagari Hindi translation protecting Tier-2/3 users.")
    ]
    for bold_txt, norm_txt in points_r:
        p = tf_r.add_paragraph()
        p.text = "[PROACTIVE] " + bold_txt + norm_txt
        p.font.size = Pt(11.5)
        p.font.color.rgb = TEXT_LIGHT

    s3.notes_slide.notes_text_frame.text = (
        "[TIME: 1:10 - 1:55]\n"
        "[WHAT TO SAY]: Current banking apps focus on authentication—verifying who is paying with a PIN. But they never inspect whether the payment context itself is fraudulent. The critical security gap exists in the 3 seconds before the PIN is entered. Suraksha AI fills this gap with a zero-trust pre-payment inspection layer. In under 200 milliseconds, we evaluate the QR payload, receipt image integrity, and message urgency to give users an explainable risk verdict before money is transferred.\n"
        "[ACTION]: Emphasize the contrast between reactive dispute forms and proactive pre-payment warnings.\n"
        "[TRANSITION]: Let's look at the core solution architecture."
    )

    # ==========================================
    # SLIDE 4: THE SOLUTION
    # ==========================================
    s4 = prs.slides.add_slide(blank_layout)
    apply_bg(s4)
    add_header(s4, "THE SOLUTION", "Meet Suraksha AI: Multi-Modal Defense Ecosystem", CYAN_ACCENT)

    pillars = [
        ("01. Live QR Scanner", "WebRTC camera feed parses UPI URI, checks VPA blacklists & validates HMAC signatures.", GREEN_SAFE),
        ("02. Image Forensic ELA", "RAM-only OpenCV JPEG Error Level Analysis detects doctored receipt screenshots.", BLUE_ACCENT),
        ("03. NLP Message Validator", "Naive Bayes & Regex heuristic matrix flags cashback / KYC urgency traps in <20ms.", AMBER_WARN),
        ("04. Signed QR Generator", "Empowers shopkeepers to generate cryptographically tamper-proof merchant QR codes.", CYAN_ACCENT)
    ]

    for i, (title, desc, col) in enumerate(pillars):
        x = 0.8 + i * 2.95
        add_card(s4, x, 1.55, 2.75, 2.35, CARD_BG, CARD_BORDER)
        
        tb = s4.shapes.add_textbox(Inches(x + 0.2), Inches(1.75), Inches(2.35), Inches(1.95))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13.5)
        p.font.bold = True
        p.font.color.rgb = col
        
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = TEXT_MUTED

    add_card(s4, 0.8, 4.15, 11.7, 2.65, RGBColor(15, 23, 42), RGBColor(56, 189, 248))
    
    tb_flow = s4.shapes.add_textbox(Inches(1.1), Inches(4.3), Inches(11.1), Inches(2.3))
    tf_f = tb_flow.text_frame
    tf_f.word_wrap = True
    
    p = tf_f.paragraphs[0]
    p.text = "THE ZERO-TRUST EVALUATION PIPELINE"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT

    p_steps = tf_f.add_paragraph()
    p_steps.text = "User Input (QR / Image / SMS)  ──>  In-Memory Sanitization  ──>  Multi-Signal Evaluation (ELA + HMAC + NLP)  ──>  Weighted Risk Engine  ──>  Actionable HUD Verdict"
    p_steps.font.size = Pt(12)
    p_steps.font.bold = True
    p_steps.font.color.rgb = TEXT_WHITE

    p_scores = tf_f.add_paragraph()
    p_scores.text = "[SAFE] Score 0–30: Trusted Merchant (Auto-Launch Payment)   |   [CAUTION] Score 31–70: Verify Identity   |   [HIGH RISK] Score 71–100: Block & Report"
    p_scores.font.size = Pt(11)
    p_scores.font.color.rgb = RGBColor(226, 232, 240)

    s4.notes_slide.notes_text_frame.text = (
        "[TIME: 1:55 - 2:15]\n"
        "[WHAT TO SAY]: Suraksha AI operates across four unified pillars: a live QR scanner, an in-memory image forensics engine, a sub-20ms NLP text validator, and a cryptographic QR generator for merchants. Our pipeline ingests the raw payload, sanitizes it in volatile RAM, runs parallel multi-signal evaluation, and produces an instant risk verdict. Let me show you this in action right now."
    )

    # ==========================================
    # SLIDE 5: LIVE PRODUCT EVIDENCE
    # ==========================================
    s5 = prs.slides.add_slide(blank_layout)
    apply_bg(s5)
    add_header(s5, "PRODUCT EVIDENCE", "See Suraksha in Action: Real-Time Risk HUD", GREEN_SAFE)

    demo_cards = [
        ("1. Multi-Modal Scanner", "realtime_scanner.png", "Live camera QR & screenshot input interface"),
        ("2. Explainable Threat HUD", "threat_breakdown_hud.png", "Calculates composite risk score with factor checklist"),
        ("3. Forensic ELA Analysis", "forensic_ela_check.png", "Exposes pixel modifications in doctored receipts")
    ]

    for i, (title, img_name, sub) in enumerate(demo_cards):
        x = 0.8 + i * 4.0
        add_card(s5, x, 1.55, 3.7, 5.25, CARD_BG, CARD_BORDER)
        
        tb = s5.shapes.add_textbox(Inches(x + 0.2), Inches(1.7), Inches(3.3), Inches(0.45))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13.5)
        p.font.bold = True
        p.font.color.rgb = CYAN_ACCENT
        
        img_path = os.path.join(assets_dir, img_name)
        if os.path.exists(img_path):
            s5.shapes.add_picture(img_path, Inches(x + 0.2), Inches(2.2), width=Inches(3.3))
            
        tb_sub = s5.shapes.add_textbox(Inches(x + 0.2), Inches(6.05), Inches(3.3), Inches(0.6))
        tf_s = tb_sub.text_frame
        tf_s.word_wrap = True
        p = tf_s.paragraphs[0]
        p.text = sub
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_MUTED

    s5.notes_slide.notes_text_frame.text = (
        "[TIME: 2:15 - 3:15]\n"
        "[WHAT TO SAY]: [LIVE DEMO DEMONSTRATION / BACKUP WALKTHROUGH]\n"
        "1. First, we input a deceptive payment message offering cashback. In milliseconds, our NLP model flags urgency patterns and intent mismatches, calculating an 85% Critical Risk Score.\n"
        "2. Second, a vendor uploads a payment screenshot. Our OpenCV Error Level Analysis inspects pixel compression noise in RAM without saving any file to disk, instantly flagging spliced text overlay.\n"
        "3. Third, with one click, the entire UI translates dynamically into Hindi, ensuring vendors across India can understand the threat clearly.\n"
        "[ACTION]: Execute live scan on localhost or present the 3 high-res UI panels.\n"
        "[TRANSITION]: Let's look at the underlying technology that powers this engine."
    )

    # ==========================================
    # SLIDE 6: TECHNOLOGY & INNOVATION
    # ==========================================
    s6 = prs.slides.add_slide(blank_layout)
    apply_bg(s6)
    add_header(s6, "TECHNICAL ARCHITECTURE", "The Intelligence Behind Suraksha AI", BLUE_ACCENT)

    tech_cols = [
        ("Layer 1: Visual Image Forensics", [
            ("75% JPEG Error Level Analysis: ", "Re-compresses in-memory byte streams to detect localized quantization noise."),
            ("Laplacian Spatial Variance: ", "Measures edge sharpness inconsistencies to detect sharp text overlays (Threshold > 2000)."),
            ("Zero-Disk Persistence: ", "Bytes processed via `io.BytesIO` in volatile RAM without saving files to disk.")
        ], CYAN_ACCENT),
        
        ("Layer 2: Cryptographic QR Shield", [
            ("HMAC SHA-256 Signatures: ", "Merchants sign VPA parameters using private keys (`sign = SHA256(pn || pa || key)`)."),
            ("Merchant Registry Lookup: ", "Matches scanned signatures against trusted databases to expose sticker swapping."),
            ("Format & Scheme Integrity: ", "Validates strictly formatted `upi://pay` URIs against spoofed deep links.")
        ], BLUE_ACCENT),
        
        ("Layer 3: NLP & Threat Intelligence", [
            ("Naive Bayes TF-IDF Model: ", "Classifies SMS/chat strings into benign alerts vs social engineering traps in <20ms."),
            ("Regex Heuristic Rules: ", "Detects high-pressure triggers (account blocked, KYC update, cashback reward)."),
            ("Bilingual Translation: ", "Regex word-boundary mapping ensures zero Hindi UI distortion.")
        ], GREEN_SAFE)
    ]

    for i, (title, bullets, col) in enumerate(tech_cols):
        x = 0.8 + i * 4.0
        add_card(s6, x, 1.55, 3.7, 5.25, CARD_BG, CARD_BORDER)
        
        tb = s6.shapes.add_textbox(Inches(x + 0.25), Inches(1.75), Inches(3.2), Inches(0.65))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13.5)
        p.font.bold = True
        p.font.color.rgb = col

        tb_c = s6.shapes.add_textbox(Inches(x + 0.25), Inches(2.55), Inches(3.2), Inches(3.9))
        tf_c = tb_c.text_frame
        tf_c.word_wrap = True
        for b_title, b_desc in bullets:
            p = tf_c.add_paragraph() if tf_c.paragraphs[0].text else tf_c.paragraphs[0]
            p.text = "• " + b_title + b_desc
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT_MUTED

    s6.notes_slide.notes_text_frame.text = (
        "[TIME: 3:15 - 3:50]\n"
        "[WHAT TO SAY]: Under the hood, Suraksha AI combines three lightweight, high-performance engines: OpenCV JPEG Error Level Analysis running in volatile RAM to protect privacy, HMAC SHA-256 cryptographic verification that validates merchant signatures to eliminate physical sticker swaps, and Multinomial Naive Bayes classification that parses conversational context in under 20 milliseconds. The entire evaluation pipeline executes with sub-200ms latency, delivering immediate protection on edge devices.\n"
        "[TRANSITION]: How does this compare to traditional market solutions?"
    )

    # ==========================================
    # SLIDE 7: DIFFERENTIATION
    # ==========================================
    s7 = prs.slides.add_slide(blank_layout)
    apply_bg(s7)
    add_header(s7, "COMPETITIVE DIFFERENTIATION", "From Payment Execution to Payment Risk Intelligence", CYAN_ACCENT)

    add_card(s7, 0.8, 1.55, 11.7, 5.25, CARD_BG, CARD_BORDER)
    
    table_shape = s7.shapes.add_table(6, 4, Inches(1.1), Inches(1.85), Inches(11.1), Inches(4.5))
    table = table_shape.table
    table.columns[0].width = Inches(3.3)
    table.columns[1].width = Inches(2.6)
    table.columns[2].width = Inches(2.6)
    table.columns[3].width = Inches(2.6)

    headers = ["Capability / Feature", "Traditional UPI Apps", "Server-Side Fraud Logs", "Suraksha AI Platform"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = CYAN_ACCENT if j == 3 else TEXT_WHITE
        p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT

    rows = [
        ("Inspection Timing", "Post-PIN Entry", "Post-Settlement Logs", "Pre-Authorization (<200ms)"),
        ("QR Sticker Swap Defense", "None (Blind trust)", "None", "HMAC SHA-256 Validation"),
        ("Screenshot Forgery Detection", "None", "None", "In-Memory ELA Forensics"),
        ("User Privacy Policy", "Server Transaction Logs", "Bulk Server Storage", "Zero-Data Retention (RAM)"),
        ("Language Inclusivity", "Static / English", "English Only", "Dynamic English & Hindi")
    ]

    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i+1, j)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            if j == 3:
                p.font.bold = True
                p.font.color.rgb = GREEN_SAFE
            elif j == 0:
                p.font.bold = True
                p.font.color.rgb = TEXT_WHITE
            else:
                p.font.color.rgb = TEXT_MUTED
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT

    s7.notes_slide.notes_text_frame.text = (
        "[TIME: 3:50 - 4:15]\n"
        "[WHAT TO SAY]: Mainstream payment apps focus on transaction throughput; they do not perform receipt forensics or cryptographic QR sticker verification. Backend banking logs catch fraud hours after settlement. Suraksha AI is built specifically for pre-authorization interception with explainable risk indicators and zero-data retention privacy.\n"
        "[TRANSITION]: Now let's explore our commercialization strategy and market scalability."
    )

    # ==========================================
    # SLIDE 8: BUSINESS + SCALE
    # ==========================================
    s8 = prs.slides.add_slide(blank_layout)
    apply_bg(s8)
    add_header(s8, "BUSINESS MODEL & SCALE", "Built for Users. Designed to Scale.", GREEN_SAFE)

    biz_pillars = [
        ("B2C: Consumer Shield", "Freemium Mobile App", [
            ("Core Scanner: ", "Free real-time QR & SMS phishing protection for everyday citizens."),
            ("Premium Security HUD: ", "Subscription tier with family protection, instant scam reporting, and cyber insurance link.")
        ], GREEN_SAFE),
        
        ("B2B: Merchant Protection", "Retail Merchant Subscriptions", [
            ("Signed QR Certification: ", "Monthly subscription for verified shopkeeper badges and tamper-proof counter stands."),
            ("POS Screenshot Verifier: ", "Instant audio chime and visual screen verification to protect shopkeepers from fake receipts.")
        ], BLUE_ACCENT),
        
        ("B2B2C / API: FinTech SDK", "Enterprise Risk-as-a-Service", [
            ("Pre-Auth Risk API: ", "Microservice API licensing for neo-banks, payment aggregators, and digital wallets."),
            ("Fraud Intelligence Feeds: ", "Shared scammer VPA blacklist telemetry for financial institutions.")
        ], CYAN_ACCENT)
    ]

    for i, (title, subtitle, bullets, col) in enumerate(biz_pillars):
        x = 0.8 + i * 4.0
        add_card(s8, x, 1.55, 3.7, 5.25, CARD_BG, CARD_BORDER)
        
        tb = s8.shapes.add_textbox(Inches(x + 0.25), Inches(1.75), Inches(3.2), Inches(0.85))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13.5)
        p.font.bold = True
        p.font.color.rgb = col
        
        p_sub = tf.add_paragraph()
        p_sub.text = subtitle
        p_sub.font.size = Pt(11)
        p_sub.font.color.rgb = TEXT_WHITE

        tb_c = s8.shapes.add_textbox(Inches(x + 0.25), Inches(2.75), Inches(3.2), Inches(3.7))
        tf_c = tb_c.text_frame
        tf_c.word_wrap = True
        for b_title, b_desc in bullets:
            p = tf_c.add_paragraph() if tf_c.paragraphs[0].text else tf_c.paragraphs[0]
            p.text = "• " + b_title + b_desc
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT_MUTED

    s8.notes_slide.notes_text_frame.text = (
        "[TIME: 4:15 - 4:40]\n"
        "[WHAT TO SAY]: Our business model scales across three layers: a freemium consumer mobile app, a ₹99/month merchant subscription providing verified tamper-proof QR stands and instant receipt verification, and enterprise B2B2C API licensing for fintech payment aggregators and digital wallets. This creates an expanding network effect where every reported scam strengthens the intelligence engine across all platforms.\n"
        "[TRANSITION]: Let's conclude with our core vision."
    )

    # ==========================================
    # SLIDE 9: THE CLOSE
    # ==========================================
    s9 = prs.slides.add_slide(blank_layout)
    apply_bg(s9)

    add_card(s9, 1.2, 0.95, 10.9, 5.6, CARD_BG, RGBColor(30, 58, 138))

    b_box = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.6), Inches(1.35), Inches(4.1), Inches(0.42))
    b_box.fill.solid()
    b_box.fill.fore_color.rgb = RGBColor(15, 23, 42)
    b_box.line.color.rgb = CYAN_ACCENT
    b_box.line.width = Pt(1)
    tf = b_box.text_frame
    p = tf.paragraphs[0]
    p.text = "THE SURAKSHA AI VISION"
    p.font.size = Pt(10.5)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT
    p.alignment = PP_ALIGN.CENTER

    tb_vis = s9.shapes.add_textbox(Inches(1.6), Inches(2.0), Inches(10.1), Inches(2.1))
    tf_v = tb_vis.text_frame
    tf_v.word_wrap = True
    p = tf_v.paragraphs[0]
    p.text = "DON'T DISCOVER THE SCAM AFTER YOU PAY.\nDISCOVER THE RISK BEFORE YOU PAY."
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.alignment = PP_ALIGN.CENTER

    p_sub = tf_v.add_paragraph()
    p_sub.text = "SURAKSHA AI — Making every digital payment an informed payment."
    p_sub.font.size = Pt(16)
    p_sub.font.color.rgb = CYAN_ACCENT
    p_sub.alignment = PP_ALIGN.CENTER

    m_boxes = [
        ("STAGE 1: PROTOTYPE", "Working PWA, Flask microservice & OpenCV ELA engine validated across attack vectors.", GREEN_SAFE),
        ("STAGE 2: VALIDATION", "Pilot trials with local retail shopkeepers in Tier-2/3 market corridors.", BLUE_ACCENT),
        ("STAGE 3: SCALE", "FinTech API licensing & multi-bank threat intelligence network.", CYAN_ACCENT)
    ]

    for i, (m_title, m_desc, col) in enumerate(m_boxes):
        x = 1.6 + i * 3.45
        card_m = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(4.45), Inches(3.2), Inches(1.5))
        card_m.fill.solid()
        card_m.fill.fore_color.rgb = RGBColor(15, 23, 42)
        card_m.line.color.rgb = col
        card_m.line.width = Pt(1)
        
        tf_m = card_m.text_frame
        tf_m.word_wrap = True
        p = tf_m.paragraphs[0]
        p.text = m_title
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = col
        
        p_desc = tf_m.add_paragraph()
        p_desc.text = m_desc
        p_desc.font.size = Pt(10)
        p_desc.font.color.rgb = TEXT_MUTED

    s9.notes_slide.notes_text_frame.text = (
        "[TIME: 4:40 - 5:00]\n"
        "[WHAT TO SAY]: Judges, digital payments should not only be fast—they must be informed and secure. Our mission is simple: Don't discover the scam after you pay. Discover the risk before you pay. We have a working prototype and a clear roadmap to secure India's digital transactions. Thank you. We are Suraksha AI.\n"
        "[ACTION]: Step forward, deliver the final one-liner with conviction, and conclude."
    )

    # Save to all target paths
    out_paths = [
        r"C:\Users\suraj\Downloads\SURAKSHA-AI.pptx",
        r"S:\Hackathon\SuRaksha\Suraksha_AI_Eureka_Pitch_2026.pptx",
        r"C:\Users\suraj\.gemini\antigravity\brain\49832c7b-902c-4026-99c5-c6c194ed2e82\SURAKSHA-AI.pptx"
    ]
    for p_out in out_paths:
        try:
            prs.save(p_out)
            print(f"Saved enhanced presentation to: {p_out}")
        except Exception as e:
            print(f"Error saving to {p_out}: {e}")

if __name__ == "__main__":
    build_enhanced_pptx()
